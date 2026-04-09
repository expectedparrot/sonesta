from __future__ import annotations

import json
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Inches, Pt

from sonesta.errors import RenderError
from sonesta.models import ChartElement, ImageElement, LoadedPresentation, ShapeElement, TableElement, TextElement, TextParagraph
from sonesta.project import load_template, load_theme
from sonesta.validation import resolve_asset_path, resolve_element_box, validate_presentation


PAGE_SIZES = {
    "standard": (10.0, 7.5),
    "widescreen": (13.333, 7.5),
}


def render_presentation(loaded: LoadedPresentation) -> dict[str, str | int | list[str]]:
    validation = validate_presentation(loaded)
    if not validation["ok"]:
        raise RenderError("presentation contains validation errors")

    prs = Presentation()
    theme = load_theme(loaded.project_root, loaded.spec.theme) if loaded.spec.theme else None
    width, height = PAGE_SIZES[(theme.page_size if theme else loaded.spec.page_size)]
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    blank_layout = prs.slide_layouts[6]
    asset_paths: set[str] = set()
    temp_asset_paths: list[Path] = []
    default_template = load_template(loaded.project_root, loaded.spec.default_template) if loaded.spec.default_template else None

    for loaded_slide in loaded.slides:
        slide = prs.slides.add_slide(blank_layout)
        if loaded_slide.spec.title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(width - 1.0), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = loaded_slide.spec.title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(24)
            title_paragraph.font.bold = True
            if theme:
                title_paragraph.font.name = theme.fonts.heading
                if "text" in theme.colors:
                    title_paragraph.font.color.rgb = _hex_to_rgb(theme.colors["text"])

        template = load_template(loaded.project_root, loaded_slide.spec.template) if loaded_slide.spec.template else default_template

        template_name = loaded_slide.spec.template or loaded.spec.default_template
        for element in sorted(
            (element for element in loaded_slide.spec.elements if element.visible),
            key=lambda item: (item.z_index or 0),
        ):
            box = resolve_element_box(element, template)
            if box is None:
                raise RenderError(f"element {element.element_id} has no resolved geometry")
            if isinstance(element, TextElement):
                _add_text(slide, element, box, theme)
            elif isinstance(element, ImageElement):
                asset = resolve_asset_path(loaded.project_root, loaded_slide.path.parent, element.path)
                asset_paths.add(str(asset))
                image_path, left, top, width_inches, height_inches = _resolve_image_placement(asset, element.fit, box)
                if image_path != asset:
                    temp_asset_paths.append(image_path)
                slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width_inches), height=Inches(height_inches))
            elif isinstance(element, ShapeElement):
                _add_shape(slide, element, box, theme)
            elif isinstance(element, TableElement):
                _add_table(slide, element, box, theme)
            elif isinstance(element, ChartElement):
                _add_chart(slide, element, box)
            else:
                raise RenderError(f"unsupported element type during render: {element.type}")

        if loaded_slide.spec.notes_path:
            notes_path = (loaded_slide.path.parent / loaded_slide.spec.notes_path).resolve()
            if notes_path.exists():
                notes_text = notes_path.read_text(encoding="utf-8")
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.text = notes_text

    output_path = resolve_output_path(loaded)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False, dir=str(output_path.parent)) as tmp:
        tmp_path = Path(tmp.name)

    prs.save(str(tmp_path))
    tmp_path.replace(output_path)

    manifest_path = output_path.with_name("manifest.json")
    manifest = {
        "presentation_id": loaded.spec.presentation_id,
        "source_path": str(loaded.presentation_path),
        "output_path": str(output_path),
        "rendered_at": datetime.now(timezone.utc).isoformat(),
        "slide_count": len(loaded.slides),
        "asset_count": len(asset_paths),
        "theme": loaded.spec.theme,
        "template_set": sorted({(slide.spec.template or loaded.spec.default_template) for slide in loaded.slides if (slide.spec.template or loaded.spec.default_template)}),
    }
    manifest_path.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    for temp_asset_path in temp_asset_paths:
        if temp_asset_path.exists():
            temp_asset_path.unlink()
    return manifest


def resolve_output_path(loaded: LoadedPresentation) -> Path:
    if loaded.spec.build.output:
        return (loaded.project_root / loaded.spec.build.output).resolve()

    default_dir = loaded.project_root / ".sonesta" / "builds" / loaded.spec.presentation_id
    return default_dir / f"{loaded.spec.presentation_id}.pptx"


def _add_text(slide, element: TextElement, box: tuple[float, float, float, float], theme) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(box[0]),
        Inches(box[1]),
        Inches(box[2]),
        Inches(box[3]),
    )
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.clear()
    line_spacing = element.line_spacing if element.line_spacing is not None else 1.3
    paragraphs = element.paragraphs or [TextParagraph(text=element.text)]
    # Default space_before by bullet level: 10pt for level 0, 3pt for level 1+
    _default_space_before = {0: 10, 1: 3}
    for index, paragraph_spec in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        text = paragraph_spec.text
        if paragraph_spec.bullet:
            text = f"{'  ' * paragraph_spec.level}• {text}"
        paragraph.text = text
        paragraph.level = paragraph_spec.level
        # Line spacing
        paragraph.line_spacing = line_spacing
        # Space before
        space_before = paragraph_spec.space_before
        if space_before is None and paragraph_spec.bullet:
            space_before = _default_space_before.get(paragraph_spec.level, 3)
        if space_before is not None:
            paragraph.space_before = Pt(space_before)
        # Space after
        if paragraph_spec.space_after is not None:
            paragraph.space_after = Pt(paragraph_spec.space_after)
        if theme:
            paragraph.font.name = theme.fonts.body
            if "text" in theme.colors:
                paragraph.font.color.rgb = _hex_to_rgb(theme.colors["text"])
            if element.style and element.style in theme.styles:
                style = theme.styles[element.style]
                if style.font_family is not None:
                    paragraph.font.name = style.font_family
                if style.font_size is not None:
                    paragraph.font.size = Pt(style.font_size)
                if style.bold is not None:
                    paragraph.font.bold = style.bold
                if style.color is not None:
                    paragraph.font.color.rgb = _hex_to_rgb(style.color)
        if element.font_size is not None:
            paragraph.font.size = Pt(element.font_size)
        if element.bold is not None:
            paragraph.font.bold = element.bold
        if paragraph_spec.bold is not None:
            paragraph.font.bold = paragraph_spec.bold
        if paragraph_spec.url is not None and paragraph.runs:
            paragraph.runs[0].hyperlink.address = paragraph_spec.url


def _add_shape(slide, element: ShapeElement, box: tuple[float, float, float, float], theme) -> None:
    if element.shape == "rect":
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(box[0]),
            Inches(box[1]),
            Inches(box[2]),
            Inches(box[3]),
        )
    elif element.shape == "ellipse":
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(box[0]),
            Inches(box[1]),
            Inches(box[2]),
            Inches(box[3]),
        )
    elif element.shape == "line":
        shape = slide.shapes.add_shape(
            MSO_SHAPE.LINE_INVERSE,
            Inches(box[0]),
            Inches(box[1]),
            Inches(box[2]),
            Inches(box[3]),
        )
    else:
        raise RenderError(f"unsupported shape kind: {element.shape}")

    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(theme.colors["accent"]) if theme and "accent" in theme.colors else RGBColor(230, 230, 230)
    shape.line.color.rgb = _hex_to_rgb(theme.colors["text"]) if theme and "text" in theme.colors else RGBColor(80, 80, 80)
    if element.text:
        shape.text_frame.text = element.text
        paragraph = shape.text_frame.paragraphs[0]
        if theme:
            paragraph.font.name = theme.fonts.body
            if "background" in theme.colors:
                paragraph.font.color.rgb = _hex_to_rgb(theme.colors["background"])


def _add_table(slide, element: TableElement, box: tuple[float, float, float, float], theme) -> None:
    rows = len(element.rows)
    cols = len(element.rows[0])
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(box[0]),
        Inches(box[1]),
        Inches(box[2]),
        Inches(box[3]),
    )
    table = table_shape.table
    if element.column_widths is not None:
        for col_index, width in enumerate(element.column_widths):
            table.columns[col_index].width = Inches(width)
    for row_index, row in enumerate(element.rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            if theme:
                paragraph.font.name = theme.fonts.body
                if "text" in theme.colors:
                    paragraph.font.color.rgb = _hex_to_rgb(theme.colors["text"])
            if row_index == 0 and element.first_row_header:
                paragraph.font.bold = True


def _add_chart(slide, element: ChartElement, box: tuple[float, float, float, float]) -> None:
    chart_type = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }[element.chart_type]
    if element.chart_type == "pie":
        data = ChartData()
        data.categories = element.categories
        data.add_series(element.series[0].name, element.series[0].values)
    else:
        data = CategoryChartData()
        data.categories = element.categories
        for series in element.series:
            data.add_series(series.name, series.values)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(box[0]),
        Inches(box[1]),
        Inches(box[2]),
        Inches(box[3]),
        data,
    )
    chart = chart_shape.chart
    chart.has_legend = element.show_legend
    if element.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = element.title
    if element.value_format:
        try:
            chart.value_axis.tick_labels.number_format = element.value_format
        except AttributeError:
            pass


def _hex_to_rgb(value: str) -> RGBColor:
    normalized = value.lstrip("#")
    return RGBColor.from_string(normalized.upper())


def _resolve_image_placement(
    asset_path: Path,
    fit: str,
    box: tuple[float, float, float, float],
) -> tuple[Path, float, float, float, float]:
    left, top, box_w, box_h = box
    if fit == "stretch":
        return asset_path, left, top, box_w, box_h

    with Image.open(asset_path) as image:
        image_w, image_h = image.size
        image_ratio = image_w / image_h
        box_ratio = box_w / box_h

        if fit == "contain":
            if image_ratio > box_ratio:
                width_inches = box_w
                height_inches = box_w / image_ratio
                return asset_path, left, top + ((box_h - height_inches) / 2), width_inches, height_inches
            height_inches = box_h
            width_inches = box_h * image_ratio
            return asset_path, left + ((box_w - width_inches) / 2), top, width_inches, height_inches

        if fit == "cover":
            if image_ratio > box_ratio:
                target_w = int(image_h * box_ratio)
                offset_x = int((image_w - target_w) / 2)
                crop_box = (offset_x, 0, offset_x + target_w, image_h)
            else:
                target_h = int(image_w / box_ratio)
                offset_y = int((image_h - target_h) / 2)
                crop_box = (0, offset_y, image_w, offset_y + target_h)
            cropped = image.crop(crop_box)
            tmp = tempfile.NamedTemporaryFile(suffix=asset_path.suffix or ".png", delete=False)
            tmp_path = Path(tmp.name)
            tmp.close()
            cropped.save(tmp_path)
            return tmp_path, left, top, box_w, box_h

    return asset_path, left, top, box_w, box_h
