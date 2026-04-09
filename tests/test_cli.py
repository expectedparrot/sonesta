from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typer.testing import CliRunner

from sonesta.cli import app
from sonesta.project import dump_json


runner = CliRunner()


def test_init_creates_project_structure(tmp_path: Path) -> None:
    result = runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["ok"] is True
    assert (tmp_path / ".sonesta" / "config.json").exists()
    assert (tmp_path / "presentations").exists()
    assert (tmp_path / "assets").exists()


def test_validate_and_render_minimal_presentation(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    presentation_dir = tmp_path / "presentations" / "demo"
    slides_dir = presentation_dir / "slides"
    slides_dir.mkdir(parents=True)

    dump_json(
        slides_dir / "001-title.json",
        {
            "slide_id": "title",
            "title": "Hello",
            "elements": [
                {
                    "element_id": "body",
                    "type": "text",
                    "x": 1,
                    "y": 1.5,
                    "w": 4,
                    "h": 1,
                    "text": "World",
                }
            ],
        },
    )
    dump_json(
        presentation_dir / "presentation.json",
        {
            "version": 1,
            "presentation_id": "demo",
            "title": "Demo",
            "slides": ["slides/001-title.json"],
            "build": {"output": ".sonesta/builds/demo/demo.pptx"},
        },
    )

    validate_result = runner.invoke(
        app,
        ["validate", str(presentation_dir / "presentation.json"), "--format", "json"],
    )
    assert validate_result.exit_code == 0
    assert json.loads(validate_result.stdout)["ok"] is True

    render_result = runner.invoke(
        app,
        ["render", str(presentation_dir / "presentation.json"), "--format", "json"],
    )
    assert render_result.exit_code == 0
    payload = json.loads(render_result.stdout)
    assert payload["ok"] is True
    assert (tmp_path / ".sonesta" / "builds" / "demo" / "demo.pptx").exists()
    assert (tmp_path / ".sonesta" / "builds" / "demo" / "manifest.json").exists()


def test_new_presentation_creates_starter_files(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    result = runner.invoke(
        app,
        ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"],
    )
    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["ok"] is True
    assert (tmp_path / "presentations" / "demo" / "presentation.json").exists()
    assert (tmp_path / "presentations" / "demo" / "slides" / "001-title.json").exists()


def test_new_slide_scaffolds_and_adds_to_presentation(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"

    result = runner.invoke(
        app,
        ["new", "slide", str(presentation_path), "roadmap", "--title", "Roadmap", "--format", "json"],
    )
    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["ok"] is True

    slide_path = tmp_path / "presentations" / "demo" / "slides" / "002-roadmap.json"
    assert slide_path.exists()
    slide = json.loads(slide_path.read_text(encoding="utf-8"))
    assert slide["slide_id"] == "roadmap"
    assert slide["elements"] == []

    listed = json.loads(runner.invoke(app, ["slides", "list", str(presentation_path), "--format", "json"]).stdout)
    assert [slide["slide_id"] for slide in listed["slides"]] == ["title", "roadmap"]


def test_slide_element_commands_build_renderable_slide(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "slide", str(tmp_path / "presentations" / "demo" / "presentation.json"), "summary", "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"
    image_path = tmp_path / "assets" / "logo.png"
    Image.new("RGB", (60, 30), color=(0, 255, 0)).save(image_path)

    result = runner.invoke(
        app,
        [
            "slides",
            "add-text",
            str(presentation_path),
            "summary",
            "headline",
            "--text",
            "Quarterly summary",
            "--x",
            "1",
            "--y",
            "1",
            "--w",
            "4",
            "--h",
            "0.6",
            "--style",
            "headline",
            "--format",
            "json",
        ],
    )
    assert result.exit_code == 0

    result = runner.invoke(
        app,
        [
            "slides",
            "add-image",
            str(presentation_path),
            "summary",
            "logo",
            "--path",
            "assets/logo.png",
            "--x",
            "8.5",
            "--y",
            "1",
            "--w",
            "2",
            "--h",
            "1",
            "--format",
            "json",
        ],
    )
    assert result.exit_code == 0

    result = runner.invoke(
        app,
        [
            "slides",
            "add-table",
            str(presentation_path),
            "summary",
            "metrics",
            "--slot",
            "hero",
            "--rows-json",
            '[["Metric","Value"],["Revenue","$10M"]]',
            "--column-widths-json",
            "[2.5,2.0]",
            "--format",
            "json",
        ],
    )
    assert result.exit_code == 0

    result = runner.invoke(
        app,
        [
            "slides",
            "add-chart",
            str(presentation_path),
            "summary",
            "trend",
            "--chart-type",
            "line",
            "--x",
            "7",
            "--y",
            "2.2",
            "--w",
            "5",
            "--h",
            "3",
            "--categories-json",
            '["Q1","Q2","Q3"]',
            "--series-json",
            '[{"name":"Sales","values":[10,12,14]}]',
            "--title",
            "Trend",
            "--format",
            "json",
        ],
    )
    assert result.exit_code == 0

    result = runner.invoke(
        app,
        [
            "slides",
            "add-shape",
            str(presentation_path),
            "summary",
            "marker",
            "--shape",
            "ellipse",
            "--x",
            "0.8",
            "--y",
            "6.0",
            "--w",
            "0.3",
            "--h",
            "0.3",
            "--format",
            "json",
        ],
    )
    assert result.exit_code == 0

    validate_result = runner.invoke(app, ["validate", str(presentation_path), "--format", "json"])
    assert validate_result.exit_code == 0
    assert json.loads(validate_result.stdout)["ok"] is True

    render_result = runner.invoke(app, ["render", str(presentation_path), "--format", "json"])
    assert render_result.exit_code == 0
    assert (tmp_path / ".sonesta" / "builds" / "demo" / "demo.pptx").exists()


def test_slide_element_commands_reject_duplicate_element_id(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"

    result = runner.invoke(
        app,
        [
            "slides",
            "add-text",
            str(presentation_path),
            "title",
            "subtitle",
            "--text",
            "duplicate",
            "--slot",
            "body",
            "--format",
            "json",
        ],
    )
    assert result.exit_code == 1
    payload = json.loads(result.stdout)
    assert payload["error"]["code"] == "validation_error"


def test_slide_element_list_update_and_remove_commands(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "slide", str(tmp_path / "presentations" / "demo" / "presentation.json"), "summary", "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"
    image_path = tmp_path / "assets" / "logo.png"
    Image.new("RGB", (80, 40), color=(0, 128, 255)).save(image_path)

    runner.invoke(
        app,
        ["slides", "add-text", str(presentation_path), "summary", "headline", "--text", "Old", "--slot", "body", "--format", "json"],
    )
    runner.invoke(
        app,
        ["slides", "add-image", str(presentation_path), "summary", "logo", "--path", "assets/logo.png", "--x", "9", "--y", "1", "--w", "2", "--h", "1", "--format", "json"],
    )
    runner.invoke(
        app,
        ["slides", "add-chart", str(presentation_path), "summary", "trend", "--chart-type", "line", "--x", "6", "--y", "2", "--w", "5", "--h", "3", "--categories-json", '["Q1","Q2"]', "--series-json", '[{"name":"Sales","values":[1,2]}]', "--format", "json"],
    )

    list_result = runner.invoke(app, ["slides", "elements", str(presentation_path), "summary", "--format", "json"])
    assert list_result.exit_code == 0
    listed = json.loads(list_result.stdout)
    assert [element["element_id"] for element in listed["elements"]] == ["headline", "logo", "trend"]

    update_text = runner.invoke(
        app,
        ["slides", "update-text", str(presentation_path), "summary", "headline", "--text", "New headline", "--style", "headline", "--format", "json"],
    )
    assert update_text.exit_code == 0

    update_image = runner.invoke(
        app,
        ["slides", "update-image", str(presentation_path), "summary", "logo", "--slot", "sidebar", "--fit", "cover", "--format", "json"],
    )
    assert update_image.exit_code == 0

    update_chart = runner.invoke(
        app,
        ["slides", "update-chart", str(presentation_path), "summary", "trend", "--title", "Updated trend", "--categories-json", '["Q1","Q2","Q3"]', "--series-json", '[{"name":"Sales","values":[1,2,3]}]', "--format", "json"],
    )
    assert update_chart.exit_code == 0

    remove_result = runner.invoke(
        app,
        ["slides", "remove-element", str(presentation_path), "summary", "logo", "--format", "json"],
    )
    assert remove_result.exit_code == 0
    removed_payload = json.loads(remove_result.stdout)
    assert removed_payload["removed"]["element_id"] == "logo"

    list_result = runner.invoke(app, ["slides", "elements", str(presentation_path), "summary", "--format", "json"])
    listed = json.loads(list_result.stdout)
    ids = [element["element_id"] for element in listed["elements"]]
    assert ids == ["headline", "trend"]
    headline = next(element for element in listed["elements"] if element["element_id"] == "headline")
    assert headline["text"] == "New headline"
    trend = next(element for element in listed["elements"] if element["element_id"] == "trend")
    assert trend["title"] == "Updated trend"

    validate_result = runner.invoke(app, ["validate", str(presentation_path), "--format", "json"])
    assert validate_result.exit_code == 0

    render_result = runner.invoke(app, ["render", str(presentation_path), "--format", "json"])
    assert render_result.exit_code == 0


def test_slide_local_inspect_validate_and_notes_commands(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "slide", str(tmp_path / "presentations" / "demo" / "presentation.json"), "summary", "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"
    slide_path = tmp_path / "presentations" / "demo" / "slides" / "002-summary.json"

    runner.invoke(
        app,
        ["slides", "add-text", str(presentation_path), "summary", "headline", "--text", "Summary", "--slot", "body", "--format", "json"],
    )

    validate_result = runner.invoke(
        app,
        ["validate", str(slide_path), "--default-template", "default", "--theme", "default", "--format", "json"],
    )
    assert validate_result.exit_code == 0
    assert json.loads(validate_result.stdout)["ok"] is True

    inspect_result = runner.invoke(
        app,
        ["inspect", str(slide_path), "--default-template", "default", "--theme", "default", "--format", "json"],
    )
    assert inspect_result.exit_code == 0
    inspect_payload = json.loads(inspect_result.stdout)
    assert inspect_payload["slide"]["slide_id"] == "summary"
    assert inspect_payload["resolved_elements"][0]["resolved_box"] == [0.8, 1.4, 11.8, 4.8]

    set_notes = runner.invoke(
        app,
        ["slides", "set-notes", str(presentation_path), "summary", "--text", "Presenter note", "--format", "json"],
    )
    assert set_notes.exit_code == 0
    notes_payload = json.loads(set_notes.stdout)
    assert notes_payload["notes_text"] == "Presenter note"

    get_notes = runner.invoke(app, ["slides", "notes", str(presentation_path), "summary", "--format", "json"])
    assert get_notes.exit_code == 0
    assert json.loads(get_notes.stdout)["notes_text"] == "Presenter note"

    slide_json = json.loads(slide_path.read_text(encoding="utf-8"))
    assert slide_json["notes_path"] == "../notes/002-summary.md"

    clear_notes = runner.invoke(app, ["slides", "clear-notes", str(presentation_path), "summary", "--format", "json"])
    assert clear_notes.exit_code == 0
    assert json.loads(runner.invoke(app, ["slides", "notes", str(presentation_path), "summary", "--format", "json"]).stdout)["notes_text"] == ""


def test_update_command_rejects_wrong_element_type(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"

    result = runner.invoke(
        app,
        ["slides", "update-image", str(presentation_path), "title", "subtitle", "--path", "assets/logo.png", "--format", "json"],
    )
    assert result.exit_code == 1
    payload = json.loads(result.stdout)
    assert payload["error"]["code"] == "validation_error"


def test_slide_rename_duplicate_remove_and_render_slide(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"

    runner.invoke(app, ["new", "slide", str(presentation_path), "summary", "--format", "json"])
    runner.invoke(app, ["slides", "set-notes", str(presentation_path), "summary", "--text", "Summary notes", "--format", "json"])

    rename_result = runner.invoke(
        app,
        ["slides", "rename", str(presentation_path), "summary", "overview", "--title", "Overview", "--format", "json"],
    )
    assert rename_result.exit_code == 0
    renamed_slide = tmp_path / "presentations" / "demo" / "slides" / "002-overview.json"
    assert renamed_slide.exists()
    presentation_json = json.loads(presentation_path.read_text(encoding="utf-8"))
    assert "slides/002-overview.json" in presentation_json["slides"]

    duplicate_result = runner.invoke(
        app,
        ["slides", "duplicate", str(presentation_path), "overview", "overview_copy", "--after", "overview", "--format", "json"],
    )
    assert duplicate_result.exit_code == 0
    duplicate_slide = tmp_path / "presentations" / "demo" / "slides" / "003-overview-copy.json"
    assert duplicate_slide.exists()
    duplicate_json = json.loads(duplicate_slide.read_text(encoding="utf-8"))
    assert duplicate_json["slide_id"] == "overview_copy"
    assert duplicate_json["notes_path"] == "../notes/003-overview-copy.md"
    assert (tmp_path / "presentations" / "demo" / "notes" / "003-overview-copy.md").exists()

    render_result = runner.invoke(
        app,
        ["render-slide", str(presentation_path), "overview_copy", "--format", "json"],
    )
    assert render_result.exit_code == 0
    assert (tmp_path / ".sonesta" / "builds" / "demo__overview_copy" / "demo__overview_copy.pptx").exists()

    remove_result = runner.invoke(
        app,
        ["slides", "remove", str(presentation_path), "overview_copy", "--delete-files", "--format", "json"],
    )
    assert remove_result.exit_code == 0
    assert not duplicate_slide.exists()
    assert not (tmp_path / "presentations" / "demo" / "notes" / "003-overview-copy.md").exists()


def test_new_slide_supports_body_preset(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"

    result = runner.invoke(
        app,
        ["new", "slide", str(presentation_path), "body_slide", "--preset", "body", "--format", "json"],
    )
    assert result.exit_code == 0
    slide = json.loads((tmp_path / "presentations" / "demo" / "slides" / "002-body-slide.json").read_text(encoding="utf-8"))
    assert slide["elements"][0]["element_id"] == "body"


def test_text_bullets_and_image_fit_render(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "slide", str(tmp_path / "presentations" / "demo" / "presentation.json"), "summary", "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"
    image_path = tmp_path / "assets" / "portrait.png"
    Image.new("RGB", (100, 200), color=(200, 10, 10)).save(image_path)

    bullets_result = runner.invoke(
        app,
        [
            "slides",
            "add-text",
            str(presentation_path),
            "summary",
            "bullets",
            "--bullets-json",
            '["One","Two","Three"]',
            "--x",
            "1",
            "--y",
            "1",
            "--w",
            "4",
            "--h",
            "2",
            "--format",
            "json",
        ],
    )
    assert bullets_result.exit_code == 0

    contain_result = runner.invoke(
        app,
        [
            "slides",
            "add-image",
            str(presentation_path),
            "summary",
            "contain_img",
            "--path",
            "assets/portrait.png",
            "--fit",
            "contain",
            "--x",
            "6",
            "--y",
            "1",
            "--w",
            "2",
            "--h",
            "2",
            "--format",
            "json",
        ],
    )
    assert contain_result.exit_code == 0

    cover_result = runner.invoke(
        app,
        [
            "slides",
            "add-image",
            str(presentation_path),
            "summary",
            "cover_img",
            "--path",
            "assets/portrait.png",
            "--fit",
            "cover",
            "--x",
            "9",
            "--y",
            "1",
            "--w",
            "2",
            "--h",
            "2",
            "--format",
            "json",
        ],
    )
    assert cover_result.exit_code == 0

    render_result = runner.invoke(app, ["render", str(presentation_path), "--format", "json"])
    assert render_result.exit_code == 0
    prs = Presentation(str(tmp_path / ".sonesta" / "builds" / "demo" / "demo.pptx"))
    slide = prs.slides[1]
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 2
    contain_shape, cover_shape = pictures
    assert contain_shape.width != cover_shape.width or contain_shape.height != cover_shape.height


def test_slides_add_list_move_remove(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    runner.invoke(app, ["new", "presentation", "demo", "--project-root", str(tmp_path), "--format", "json"])
    presentation_path = tmp_path / "presentations" / "demo" / "presentation.json"
    extra_slide = tmp_path / "presentations" / "demo" / "slides" / "002-extra.json"
    dump_json(
        extra_slide,
        {
            "slide_id": "extra",
            "title": "Extra",
            "elements": [{"element_id": "body", "type": "text", "slot": "body", "text": "More"}],
        },
    )

    add_result = runner.invoke(
        app,
        ["slides", "add", str(presentation_path), "--slide", str(extra_slide), "--after", "title", "--format", "json"],
    )
    assert add_result.exit_code == 0

    list_result = runner.invoke(app, ["slides", "list", str(presentation_path), "--format", "json"])
    listed = json.loads(list_result.stdout)
    assert [slide["slide_id"] for slide in listed["slides"]] == ["title", "extra"]

    move_result = runner.invoke(
        app,
        ["slides", "move", str(presentation_path), "title", "--after", "extra", "--format", "json"],
    )
    assert move_result.exit_code == 0

    list_result = runner.invoke(app, ["slides", "list", str(presentation_path), "--format", "json"])
    listed = json.loads(list_result.stdout)
    assert [slide["slide_id"] for slide in listed["slides"]] == ["extra", "title"]

    remove_result = runner.invoke(
        app,
        ["slides", "remove", str(presentation_path), "extra", "--format", "json"],
    )
    assert remove_result.exit_code == 0
    listed = json.loads(runner.invoke(app, ["slides", "list", str(presentation_path), "--format", "json"]).stdout)
    assert [slide["slide_id"] for slide in listed["slides"]] == ["title"]


def test_validate_reports_missing_image_asset(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    presentation_dir = tmp_path / "presentations" / "demo"
    slides_dir = presentation_dir / "slides"
    slides_dir.mkdir(parents=True)

    dump_json(
        slides_dir / "001-title.json",
        {
            "slide_id": "title",
            "title": "Hello",
            "elements": [
                {
                    "element_id": "hero",
                    "type": "image",
                    "x": 1,
                    "y": 1,
                    "w": 4,
                    "h": 3,
                    "path": "assets/missing.png",
                }
            ],
        },
    )
    dump_json(
        presentation_dir / "presentation.json",
        {
            "version": 1,
            "presentation_id": "demo",
            "slides": ["slides/001-title.json"],
        },
    )

    result = runner.invoke(
        app,
        ["validate", str(presentation_dir / "presentation.json"), "--format", "json"],
    )
    assert result.exit_code == 1
    payload = json.loads(result.stdout)
    assert payload["ok"] is False
    assert payload["issues"][0]["code"] == "missing_asset"


def test_assets_inspect_reports_image_metadata(tmp_path: Path) -> None:
    asset_path = tmp_path / "sample.png"
    image = Image.new("RGB", (120, 80), color=(255, 0, 0))
    image.save(asset_path)

    result = runner.invoke(app, ["assets", "inspect", str(asset_path), "--format", "json"])
    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["ok"] is True
    assert payload["image"]["width_px"] == 120
    assert payload["image"]["height_px"] == 80
    assert payload["suffix"] == ".png"


def test_list_commands_report_project_resources(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    image_path = tmp_path / "assets" / "logo.png"
    image = Image.new("RGB", (40, 20), color=(0, 0, 255))
    image.save(image_path)

    themes_result = runner.invoke(app, ["themes", "--project-root", str(tmp_path), "--format", "json"])
    templates_result = runner.invoke(app, ["templates", "--project-root", str(tmp_path), "--format", "json"])
    assets_result = runner.invoke(app, ["assets", "list", "--project-root", str(tmp_path), "--format", "json"])

    assert themes_result.exit_code == 0
    assert templates_result.exit_code == 0
    assert assets_result.exit_code == 0

    assert json.loads(themes_result.stdout)["themes"] == ["default"]
    assert json.loads(templates_result.stdout)["templates"] == ["default"]
    assets_payload = json.loads(assets_result.stdout)
    assert assets_payload["assets"][0]["relative_path"] == "logo.png"
    assert assets_payload["assets"][0]["suffix"] == ".png"


def test_render_supports_template_slots_table_and_chart(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    presentation_dir = tmp_path / "presentations" / "demo"
    slides_dir = presentation_dir / "slides"
    slides_dir.mkdir(parents=True)

    dump_json(
        slides_dir / "001-summary.json",
        {
            "slide_id": "summary",
            "title": "Summary",
            "template": "default",
            "elements": [
                {
                    "element_id": "table_1",
                    "type": "table",
                    "slot": "hero",
                    "rows": [["Metric", "Value"], ["Revenue", "$10M"], ["Growth", "24%"]],
                },
                {
                    "element_id": "chart_1",
                    "type": "chart",
                    "x": 7.1,
                    "y": 1.4,
                    "w": 5.3,
                    "h": 3.6,
                    "chart_type": "column",
                    "title": "Quarterly Sales",
                    "show_legend": True,
                    "value_format": "0",
                    "categories": ["Q1", "Q2", "Q3"],
                    "series": [{"name": "Sales", "values": [10, 12, 14]}],
                },
            ],
        },
    )
    dump_json(
        presentation_dir / "presentation.json",
        {
            "version": 1,
            "presentation_id": "demo",
            "title": "Demo",
            "theme": "default",
            "default_template": "default",
            "slides": ["slides/001-summary.json"],
        },
    )

    validate_result = runner.invoke(app, ["validate", str(presentation_dir / "presentation.json"), "--format", "json"])
    assert validate_result.exit_code == 0
    assert json.loads(validate_result.stdout)["ok"] is True

    inspect_result = runner.invoke(app, ["inspect", str(presentation_dir / "presentation.json"), "--format", "json"])
    inspect_payload = json.loads(inspect_result.stdout)
    assert inspect_payload["slides"][0]["resolved_elements"][0]["resolved_box"] == [0.8, 1.4, 5.8, 4.8]

    render_result = runner.invoke(app, ["render", str(presentation_dir / "presentation.json"), "--format", "json"])
    assert render_result.exit_code == 0
    assert (tmp_path / ".sonesta" / "builds" / "demo" / "demo.pptx").exists()


def test_validate_rejects_bad_table_column_widths(tmp_path: Path) -> None:
    runner.invoke(app, ["init", str(tmp_path), "--format", "json"])
    presentation_dir = tmp_path / "presentations" / "demo"
    slides_dir = presentation_dir / "slides"
    slides_dir.mkdir(parents=True)

    dump_json(
        slides_dir / "001-summary.json",
        {
            "slide_id": "summary",
            "elements": [
                {
                    "element_id": "table_1",
                    "type": "table",
                    "x": 1,
                    "y": 1,
                    "w": 5,
                    "h": 3,
                    "rows": [["A", "B"], ["1", "2"]],
                    "column_widths": [2.0],
                }
            ],
        },
    )
    dump_json(
        presentation_dir / "presentation.json",
        {"version": 1, "presentation_id": "demo", "slides": ["slides/001-summary.json"]},
    )

    result = runner.invoke(app, ["validate", str(presentation_dir / "presentation.json"), "--format", "json"])
    assert result.exit_code == 1
    payload = json.loads(result.stdout)
    assert payload["error"]["code"] == "schema_error"
